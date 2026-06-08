"""Add valuation addendum slides to the AISAN PPTX deliverable.

Run from repo root:
    PYTHONPATH=/tmp/aisan_pydeps python3 -m src.report.update_pptx_addendum

The script uses python-pptx for shape creation and a small internal slide-order
adjustment so the source appendix remains the final slide.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
PPTX = ROOT / "docs" / "AISAN_4667_Take_Private_Case_Study.pptx"

NAVY = RGBColor(20, 38, 58)
BLUE = RGBColor(31, 78, 121)
TEAL = RGBColor(0, 112, 112)
GREEN = RGBColor(62, 123, 80)
AMBER = RGBColor(181, 122, 20)
RED = RGBColor(168, 55, 55)
GREY = RGBColor(95, 105, 112)
LIGHT_GREY = RGBColor(239, 242, 245)
MID_GREY = RGBColor(216, 222, 228)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(30, 34, 38)


def clear_slide(slide):
    for shape in list(slide.shapes):
        slide.shapes._spTree.remove(shape._element)


def add_text(slide, x, y, w, h, text, size=12, color=BLACK, bold=False,
             align=PP_ALIGN.LEFT, fill=None, margin=0.08, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    if fill is not None:
        box.fill.solid()
        box.fill.fore_color.rgb = fill
        box.line.color.rgb = fill
    return box


def add_title(slide, section, title, subtitle):
    add_text(slide, 0.35, 0.25, 0.55, 0.34, section, 13, WHITE, True,
             PP_ALIGN.CENTER, NAVY, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, 1.00, 0.19, 5.2, 0.36, title, 19, NAVY, True)
    add_text(slide, 1.00, 0.57, 11.2, 0.35, subtitle, 9.2, GREY)
    line = slide.shapes.add_shape(1, Inches(0.35), Inches(0.95), Inches(12.55), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = MID_GREY
    line.line.color.rgb = MID_GREY


def add_footer(slide, page_no):
    add_text(slide, 0.35, 7.08, 5.0, 0.22, "Project Measure  ·  Strictly Private & Confidential",
             7.5, GREY)
    add_text(slide, 12.45, 7.08, 0.45, 0.22, str(page_no), 7.5, GREY, align=PP_ALIGN.RIGHT)


def add_panel_title(slide, x, y, w, title, accent=BLUE):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(0.32))
    shape.fill.solid()
    shape.fill.fore_color.rgb = accent
    shape.line.color.rgb = accent
    add_text(slide, x + 0.04, y + 0.035, w - 0.08, 0.22, title, 8.5, WHITE, True,
             PP_ALIGN.LEFT, None, margin=0.02, valign=MSO_ANCHOR.MIDDLE)


def add_table(slide, x, y, w, h, headers, rows, widths=None, header_color=BLUE,
              font_size=7.0):
    table_shape = slide.shapes.add_table(
        len(rows) + 1, len(headers), Inches(x), Inches(y), Inches(w), Inches(h)
    )
    table = table_shape.table
    if widths:
        for col, width in enumerate(widths):
            table.columns[col].width = Inches(width)
    for col, header in enumerate(headers):
        cell = table.cell(0, col)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                r.font.name = "Aptos"
                r.font.size = Pt(font_size)
                r.font.bold = True
                r.font.color.rgb = WHITE
    for r_idx, row in enumerate(rows, start=1):
        for c_idx, value in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = str(value)
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_GREY if r_idx % 2 else WHITE
            cell.margin_left = Inches(0.04)
            cell.margin_right = Inches(0.04)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT if c_idx == 0 else PP_ALIGN.CENTER
                for run in p.runs:
                    run.font.name = "Aptos"
                    run.font.size = Pt(font_size)
                    run.font.color.rgb = BLACK
    return table_shape


def add_callout(slide, x, y, w, h, label, value, sub, color=BLUE):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = MID_GREY
    add_text(slide, x + 0.08, y + 0.08, w - 0.16, 0.20, label, 7.0, GREY, True,
             PP_ALIGN.CENTER)
    add_text(slide, x + 0.08, y + 0.34, w - 0.16, 0.34, value, 15.5, color, True,
             PP_ALIGN.CENTER)
    add_text(slide, x + 0.08, y + 0.76, w - 0.16, 0.28, sub, 7.0, GREY,
             align=PP_ALIGN.CENTER)


def add_valuation_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    clear_slide(slide)
    add_title(
        slide,
        "A1",
        "Valuation Context",
        "Strategic scarcity supports a premium, but sponsor returns do not clear at full precedent-style premia.",
    )

    add_panel_title(slide, 0.55, 1.20, 5.80, "AISAN entry valuation and peer frame", BLUE)
    add_callout(slide, 0.70, 1.70, 1.60, 1.05, "ENTRY EV / EBITDA", "10.1x", "FY26E EBITDA", BLUE)
    add_callout(slide, 2.50, 1.70, 1.60, 1.05, "ENTRY EV / EBIT", "13.4x", "FY26E EBIT", TEAL)
    add_callout(slide, 4.30, 1.70, 1.60, 1.05, "OFFER P / E", "31.8x", "ex-treasury EPS", AMBER)

    add_table(
        slide,
        0.65,
        3.10,
        5.55,
        1.90,
        ["Reference set", "Why relevant", "Limitation"],
        [
            ["Fukui / Zenrin", "JP software / map data", "Different mix and scale"],
            ["Topcon / Trimble", "Positioning / survey stack", "Hardware-heavy profile"],
            ["Autodesk / Bentley", "Construction software", "Much higher quality / scale"],
            ["TomTom / HD maps", "Mobility / map exposure", "Different end-market maturity"],
        ],
        widths=[1.35, 2.00, 2.20],
        header_color=BLUE,
        font_size=6.9,
    )
    add_text(
        slide,
        0.65,
        5.18,
        5.55,
        0.50,
        "Live peer EV/EBITDA was not refreshed in this pass; use a market-data terminal or verified API before bid committee.",
        7.2,
        GREY,
    )

    add_panel_title(slide, 6.65, 1.20, 5.85, "Japan take-private precedent context", TEAL)
    add_table(
        slide,
        6.75,
        1.70,
        5.65,
        3.05,
        ["Target", "Buyer", "Premium", "Relevance"],
        [
            ["Tecnos Japan", "Ant Capital", "+39%", "Small-cap IT / software"],
            ["Kaonavi", "Carlyle", "+121%", "Growth SaaS buyout"],
            ["Topcon", "KKR-led", "+17% / +58%", "Survey / positioning adjacencies"],
            ["PASCO", "SECOM / ITOCHU", "Verify", "Closest JP geospatial precedent"],
        ],
        widths=[1.20, 1.25, 1.15, 2.05],
        header_color=TEAL,
        font_size=6.8,
    )
    add_text(
        slide,
        6.75,
        4.95,
        5.65,
        0.55,
        "Reported precedent premia commonly start above 30%, but AISAN's 20% hurdle walk-to-price is only approx. JPY2,031/share.",
        7.6,
        GREY,
    )

    shape = slide.shapes.add_shape(1, Inches(0.55), Inches(5.88), Inches(11.95), Inches(0.62))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(250, 247, 240)
    shape.line.color.rgb = RGBColor(230, 210, 170)
    add_text(
        slide,
        0.72,
        6.01,
        11.55,
        0.31,
        "Takeaway: AISAN may be strategically attractive, but a financial sponsor needs price discipline, governance clearance and DD-backed value-up rather than relying on market-premium comparables.",
        8.4,
        BLACK,
        True,
        PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        0.55,
        6.60,
        11.95,
        0.26,
        "Sources: FACT_CHECK_2026-06.md, local model, reported public transaction summaries; all transaction multiples require final source-file verification.",
        6.4,
        GREY,
    )
    add_footer(slide, 17)
    return slide


def add_walk_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    clear_slide(slide)
    add_title(
        slide,
        "A2",
        "Walk-to-Price & Leverage",
        "Treasury-share correction improves returns, but base case remains below a clean 20% PE hurdle.",
    )

    add_panel_title(slide, 0.55, 1.20, 5.85, "20% hurdle walk-to-price", BLUE)
    add_callout(slide, 0.75, 1.68, 1.55, 1.05, "MAX OFFER", "JPY2,031", "20% gross IRR", BLUE)
    add_callout(slide, 2.48, 1.68, 1.55, 1.05, "VS SPOT", "+20.9%", "JPY1,680 ref", GREEN)
    add_callout(slide, 4.21, 1.68, 1.55, 1.05, "VS 6M VWAP", "+5.3%", "JPY1,929", AMBER)
    add_table(
        slide,
        0.70,
        3.05,
        5.50,
        1.70,
        ["Bid metric", "Current case", "20% hurdle cap"],
        [
            ["Offer / share", "JPY2,302", "JPY2,031"],
            ["Premium vs spot", "+37.0%", "+20.9%"],
            ["Premium vs 6m VWAP", "+19.3%", "+5.3%"],
            ["Base gross IRR", "15.6%", "20.0%"],
        ],
        widths=[2.00, 1.75, 1.75],
        header_color=BLUE,
        font_size=7.2,
    )
    add_text(
        slide,
        0.70,
        4.95,
        5.50,
        0.55,
        "Implication: do not bid above the walk-to-price unless investigation findings are clean and incremental value-up is explicitly underwritten.",
        7.8,
        GREY,
    )

    add_panel_title(slide, 6.70, 1.20, 5.80, "Illustrative leverage sensitivity", TEAL)
    add_table(
        slide,
        6.82,
        1.70,
        5.55,
        2.58,
        ["Debt / EBITDA", "Debt", "Sponsor eq.", "MOIC", "IRR"],
        [
            ["1.0x", "JPY0.8bn", "JPY8.5bn", "2.07x", "15.6%"],
            ["2.0x", "JPY1.6bn", "JPY7.7bn", "2.17x", "16.8%"],
            ["2.5x", "JPY2.0bn", "JPY7.3bn", "2.23x", "17.4%"],
            ["3.0x", "JPY2.4bn", "JPY6.9bn", "2.30x", "18.1%"],
        ],
        widths=[1.30, 1.05, 1.15, 0.95, 1.10],
        header_color=TEAL,
        font_size=7.2,
    )
    add_text(
        slide,
        6.82,
        4.55,
        5.55,
        0.64,
        "Even 3.0x gross debt does not clear 20% and would be imprudent before QoE, cash seasonality, investigation impact and lender appetite are cleared.",
        7.8,
        GREY,
    )

    shape = slide.shapes.add_shape(1, Inches(0.55), Inches(5.82), Inches(11.95), Inches(0.72))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(241, 247, 246)
    shape.line.color.rgb = RGBColor(180, 215, 210)
    add_text(
        slide,
        0.75,
        5.98,
        11.55,
        0.35,
        "Takeaway: this is not a classic leverage-driven LBO. Returns are driven by entry price, cash availability, EBITDA growth, margin expansion and exit multiple.",
        8.6,
        BLACK,
        True,
        PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        0.55,
        6.62,
        11.95,
        0.25,
        "Sources: Excel model Returns and Lev_Sensitivity tabs; local market snapshot as of 5-Jun-2026.",
        6.4,
        GREY,
    )
    add_footer(slide, 18)
    return slide


def update_source_slide(source_slide):
    for shape in source_slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        text = shape.text_frame.text.strip()
        if text == "A":
            shape.text_frame.text = "A3"
        elif text == "17":
            shape.text_frame.text = "19"
        elif text.startswith("Reported / estimated / to be verified"):
            shape.text_frame.text = (
                "Reported / estimated / to be verified in DD: founder / related-holder stake, "
                "segment margins, peer multiples, exact precedent premium bases, investment-securities "
                "mark, FY2026E guidance and Akisoku impact. See FACT_CHECK_2026-06.md."
            )


def move_original_source_to_end(prs):
    sld_ids = prs.slides._sldIdLst
    original_source = sld_ids[16]
    sld_ids.remove(original_source)
    sld_ids.append(original_source)


def main():
    prs = Presentation(PPTX)
    if len(prs.slides) != 17:
        raise RuntimeError(f"Expected 17 slides before addendum, found {len(prs.slides)}")

    source_slide = prs.slides[16]
    add_valuation_slide(prs)
    add_walk_slide(prs)
    update_source_slide(source_slide)
    move_original_source_to_end(prs)
    prs.save(PPTX)
    print(f"updated {PPTX} to {len(prs.slides)} slides")


if __name__ == "__main__":
    main()
